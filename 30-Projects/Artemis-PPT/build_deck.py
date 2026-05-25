"""
Artemis II 발표 PPT 빌드 스크립트 (3분 영어 발표용)
- python-pptx + Pillow 사용
- Wikimedia Special:FilePath 로 이미지 자동 다운로드 (실패 시 단색 폴백)

실행:
    python build_deck.py
출력:
    artemis_mission.pptx  (스크립트와 같은 폴더)

필요 패키지: python-pptx, Pillow
    (msys2) pacman -S mingw-w64-ucrt-x86_64-python-pptx
    (pip)   pip install python-pptx Pillow
"""
import os, ssl, tempfile, urllib.request, urllib.parse
from PIL import Image, ImageEnhance
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

HERE = os.path.dirname(os.path.abspath(__file__))
OUT = os.path.join(HERE, "artemis_mission.pptx")
TMP = tempfile.mkdtemp(prefix="artemis_img_")
UA = "Mozilla/5.0 (Windows NT 10.0; Win64; x64) ArtemisPPT/1.0"

ctx = ssl.create_default_context()
ctx.check_hostname = False
ctx.verify_mode = ssl.CERT_NONE

# ---- palette ----
DARK  = RGBColor(0x0A, 0x16, 0x28)
NAVY  = RGBColor(0x0D, 0x1F, 0x3C)
ORANGE= RGBColor(0xE8, 0x67, 0x0A)
BLUE  = RGBColor(0x4A, 0x90, 0xD9)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LGRAY = RGBColor(0x8B, 0xA7, 0xC7)

HEAD = "Georgia"
BODY = "Calibri"

# ---- images: Wikimedia filenames (Special:FilePath auto-resolves the hash) ----
IMAGES = {
    "moon":      "FullMoon2010.jpg",
    "earth":     "The Earth seen from Apollo 17.jpg",
    "aldrin":    "Aldrin Apollo 11 original.jpg",
    "mars":      "OSIRIS Mars true color.jpg",
    "sun":       "The Sun by the Atmospheric Imaging Assembly of NASA's Solar Dynamics Observatory - 20100819.jpg",
    "artemis":   "Artemis I launch from LC-39B (NHQ202211160029).jpg",
    "southpole": "Lunar south pole Clementine color mosaic.jpg",
    "lunarbase": "Lunar base concept drawing.jpg",
}
_cache = {}

def fetch(key, width=1400):
    """Download an image by key; return local path or None on failure."""
    if key in _cache:
        return _cache[key]
    fname = IMAGES[key]
    url = "https://commons.wikimedia.org/wiki/Special:FilePath/" + urllib.parse.quote(fname) + f"?width={width}"
    dest = os.path.join(TMP, key + ".img")
    try:
        req = urllib.request.Request(url, headers={"User-Agent": UA})
        with urllib.request.urlopen(req, timeout=30, context=ctx) as r:
            data = r.read()
        with open(dest, "wb") as f:
            f.write(data)
        Image.open(dest).verify()  # sanity check
        _cache[key] = dest
        print(f"  [ok]   {key}: {len(data)} bytes")
        return dest
    except Exception as e:
        print(f"  [FAIL] {key}: {type(e).__name__} {str(e)[:50]}")
        _cache[key] = None
        return None

def cover(path, w_in, h_in, darken=0.0, dpi=200):
    """Crop-to-cover (w_in x h_in), optionally darken (0..1). Returns processed png path."""
    img = Image.open(path).convert("RGB")
    tw, th = int(w_in * dpi), int(h_in * dpi)
    sw, sh = img.size
    scale = max(tw / sw, th / sh)
    nw, nh = int(sw * scale), int(sh * scale)
    img = img.resize((nw, nh), Image.LANCZOS)
    left, top = (nw - tw) // 2, (nh - th) // 2
    img = img.crop((left, top, left + tw, top + th))
    if darken > 0:
        img = ImageEnhance.Brightness(img).enhance(1.0 - darken)
    out = os.path.join(TMP, f"c_{os.path.basename(path)}_{tw}x{th}_{int(darken*100)}.png")
    img.save(out)
    return out

# ---- pptx helpers ----
prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625)
BLANK = prs.slide_layouts[6]

def new_slide(bg=DARK):
    s = prs.slides.add_slide(BLANK)
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = bg
    return s

def rect(s, x, y, w, h, color, line_color=None, line_w=1.0):
    from pptx.enum.shapes import MSO_SHAPE
    sp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    sp.fill.solid(); sp.fill.fore_color.rgb = color
    if line_color is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line_color; sp.line.width = Pt(line_w)
    sp.shadow.inherit = False
    return sp

def picture(s, key, x, y, w, h, darken=0.0, fallback=NAVY):
    p = fetch(key)
    if p is None:
        rect(s, x, y, w, h, fallback)
        return
    cp = cover(p, w, h, darken=darken)
    s.shapes.add_picture(cp, Inches(x), Inches(y), Inches(w), Inches(h))

def text(s, x, y, w, h, lines, size=14, color=WHITE, bold=False, font=BODY,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, spacing=None):
    """lines: str or list[str]. Each list item = one paragraph."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True
    tf.margin_left = 0; tf.margin_right = 0; tf.margin_top = 0; tf.margin_bottom = 0
    tf.vertical_anchor = anchor
    if isinstance(lines, str):
        lines = [lines]
    for i, ln in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = ln
        r.font.size = Pt(size); r.font.bold = bold; r.font.name = font
        r.font.color.rgb = color
        if spacing is not None:
            from pptx.util import Pt as _Pt
            p.space_after = _Pt(spacing)
    return tb

# =================================================================
print("Downloading images + building slides...")

# ---- SLIDE 1: TITLE ----
s = new_slide(DARK)
picture(s, "moon", 5.8, 0, 4.2, 5.625, darken=0.30, fallback=NAVY)
rect(s, 0.5, 1.25, 0.07, 2.6, ORANGE)
text(s, 0.7, 1.15, 5.0, 2.6, ["WHY ARE WE GOING", "BACK TO THE MOON?"],
     size=36, color=WHITE, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.7, 4.0, 5.0, 0.5, "Artemis II Mission", size=18, color=ORANGE, bold=True)
text(s, 0.7, 4.6, 5.0, 0.4, "Humanity's Return to Deep Space", size=13, color=LGRAY)

# ---- SLIDE 2: INTRODUCTION ----
s = new_slide(DARK)
picture(s, "artemis", 5.4, 0, 4.6, 5.625, darken=0.35, fallback=NAVY)
text(s, 0.5, 0.4, 4.5, 0.3, "INTRODUCTION", size=10, color=ORANGE, bold=True)
text(s, 0.5, 0.75, 4.6, 1.8, ["50 Years Later,", "We Return"],
     size=32, color=WHITE, bold=True, font=HEAD)
text(s, 0.5, 2.7, 4.3, 0.9,
     "NASA's Artemis II recently broke the deep-space distance record set by Apollo - over 50 years ago.",
     size=14, color=LGRAY)
rect(s, 0.5, 3.85, 4.3, 1.0, ORANGE)
text(s, 0.6, 3.85, 4.1, 1.0, "So... why are we spending billions to go back?",
     size=14, color=WHITE, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)

# ---- SLIDE 3: RESOURCES ----
s = new_slide(DARK)
text(s, 0.5, 0.22, 9, 0.28, "REASON 1", size=10, color=ORANGE, bold=True)
text(s, 0.5, 0.5, 9, 0.75, "The Moon Is a Treasure Chest", size=28, color=WHITE, bold=True, font=HEAD)
rect(s, 0.5, 1.28, 9, 0.04, ORANGE)
# card 1
rect(s, 0.4, 1.45, 4.5, 3.7, NAVY)
picture(s, "southpole", 0.4, 1.45, 4.5, 2.0, darken=0.0, fallback=BLUE)
rect(s, 0.4, 2.95, 1.5, 0.35, ORANGE)
text(s, 0.4, 2.95, 1.5, 0.35, "WATER ICE", size=9, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 0.62, 3.4, 4.1, 0.4, "South Pole Deposits", size=15, color=WHITE, bold=True, font=HEAD)
text(s, 0.62, 3.85, 4.0, 1.2,
     "Water ice at the lunar South Pole can be turned into rocket fuel - directly in space.",
     size=12, color=LGRAY)
# card 2
rect(s, 5.1, 1.45, 4.5, 3.7, NAVY)
picture(s, "sun", 5.1, 1.45, 4.5, 2.0, darken=0.0, fallback=ORANGE)
rect(s, 5.1, 2.95, 1.6, 0.35, BLUE)
text(s, 5.1, 2.95, 1.6, 0.35, "HELIUM-3", size=9, color=WHITE, bold=True,
     align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
text(s, 5.32, 3.4, 4.1, 0.4, "Future Clean Energy", size=15, color=WHITE, bold=True, font=HEAD)
text(s, 5.32, 3.85, 4.0, 1.2,
     "The Moon holds large amounts of Helium-3, a rare material that may power a clean-energy future.",
     size=12, color=LGRAY)

# ---- SLIDE 4: SURVIVAL ----
s = new_slide(DARK)
picture(s, "aldrin", 5.4, 0, 4.6, 5.625, darken=0.35, fallback=NAVY)
text(s, 0.5, 0.4, 4.5, 0.28, "REASON 2", size=10, color=ORANGE, bold=True)
text(s, 0.5, 0.72, 4.6, 1.8, ["A Survival", "Training Ground"],
     size=30, color=WHITE, bold=True, font=HEAD)
text(s, 0.5, 2.7, 4.3, 0.6,
     "Before going to Mars, we must practice long-term survival in space.",
     size=13, color=LGRAY)
items = [
    ("Life-Support Systems", "Test the equipment that keeps humans alive in space"),
    ("Radiation Protection", "Learn to shield astronauts from dangerous space radiation"),
    ("Long-Term Habitats", "Design living spaces built to last in deep space"),
]
for i, (t, d) in enumerate(items):
    y = 3.4 + i * 0.72
    rect(s, 0.5, y + 0.06, 0.09, 0.09, ORANGE)
    text(s, 0.75, y, 3.9, 0.28, t, size=13, color=WHITE, bold=True)
    text(s, 0.75, y + 0.3, 3.9, 0.3, d, size=11, color=LGRAY)

# ---- SLIDE 5: LUNAR BASE ----
s = new_slide(DARK)
picture(s, "lunarbase", 3.6, 0, 6.4, 5.625, darken=0.30, fallback=NAVY)
rect(s, 0, 0, 4.0, 5.625, DARK)
text(s, 0.5, 0.4, 3.2, 0.28, "REASON 3", size=10, color=ORANGE, bold=True)
text(s, 0.5, 0.72, 3.4, 1.8, ["Building a Base", "on the Moon"],
     size=27, color=WHITE, bold=True, font=HEAD)
rect(s, 0.5, 2.6, 2.8, 0.04, ORANGE)
text(s, 0.5, 2.8, 3.2, 1.1,
     "NASA plans a permanent lunar base where astronauts can live and work on the Moon's surface.",
     size=13, color=LGRAY)
rect(s, 0.5, 4.05, 3.2, 1.15, NAVY, line_color=BLUE, line_w=1.0)
text(s, 0.65, 4.05, 3.0, 1.15, "A launching point for the first human mission to Mars.",
     size=13, color=WHITE, bold=True, font=HEAD, anchor=MSO_ANCHOR.MIDDLE)

# ---- SLIDE 6: CONCLUSION ----
s = new_slide(DARK)
picture(s, "earth", 6.0, 0, 4.0, 5.625, darken=0.30, fallback=NAVY)
text(s, 0.5, 0.4, 5.5, 0.28, "CONCLUSION", size=10, color=ORANGE, bold=True)
text(s, 0.5, 0.72, 5.5, 1.8, ["Not a Waste -", "An Investment"],
     size=34, color=WHITE, bold=True, font=HEAD)
cards = [("01", "Resources", "Water ice & Helium-3"),
         ("02", "Survival Tech", "Life support & radiation"),
         ("03", "Path to Mars", "Lunar base as launchpad")]
for i, (n, t, d) in enumerate(cards):
    x = 0.5 + i * 3.1
    rect(s, x, 2.85, 2.8, 1.5, NAVY)
    text(s, x + 0.15, 2.9, 0.9, 0.5, n, size=20, color=ORANGE, bold=True, font=HEAD)
    text(s, x + 0.15, 3.42, 2.5, 0.3, t, size=13, color=WHITE, bold=True)
    text(s, x + 0.15, 3.78, 2.5, 0.5, d, size=11, color=LGRAY)
rect(s, 0.5, 4.6, 9.0, 0.7, ORANGE)
text(s, 0.6, 4.6, 8.8, 0.7,
     "This return to the Moon may be humanity's first step toward becoming a multi-planetary species.",
     size=13, color=WHITE, bold=True, font=HEAD, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

prs.save(OUT)
print(f"\nDone -> {OUT}")
